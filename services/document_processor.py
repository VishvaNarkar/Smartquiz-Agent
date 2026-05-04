from __future__ import annotations

import io
import os
import re
import tempfile
from pathlib import Path


SUPPORTED_DOCUMENT_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".doc", ".docx"}
MAX_DOCUMENT_BYTES = 12 * 1024 * 1024
MAX_EXTRACTED_CHARACTERS = 12000


def get_document_extension(filename: str) -> str:
    return Path((filename or "").strip().lower()).suffix


def is_supported_document(filename: str) -> bool:
    return get_document_extension(filename) in SUPPORTED_DOCUMENT_EXTENSIONS


def _normalize_text(text: str) -> str:
    cleaned = re.sub(r"\s+", " ", (text or "").strip())
    if len(cleaned) > MAX_EXTRACTED_CHARACTERS:
        cleaned = cleaned[:MAX_EXTRACTED_CHARACTERS].rstrip()
    return cleaned


def _extract_pdf_text(contents: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("PDF support is unavailable. Install pypdf to enable PDF uploads.") from exc

    reader = PdfReader(io.BytesIO(contents))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx_text(contents: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("DOCX support is unavailable. Install python-docx to enable DOCX uploads.") from exc

    document = Document(io.BytesIO(contents))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)
    return "\n".join(paragraphs)


def _extract_pptx_text(contents: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX support is unavailable. Install python-pptx to enable slide uploads.") from exc

    presentation = Presentation(io.BytesIO(contents))
    text_chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_chunks.append(shape.text)
    return "\n".join(text_chunks)


def _extract_legacy_doc_text(contents: bytes) -> str:
    try:
        import textract  # type: ignore
    except ImportError as exc:
        raise RuntimeError(
            "Legacy .doc files are not supported in this environment. Please convert the file to .docx or PDF."
        ) from exc

    with tempfile.NamedTemporaryFile(delete=False, suffix=".doc") as tmp:
        tmp.write(contents)
        temp_path = tmp.name

    try:
        extracted = textract.process(temp_path)
        return extracted.decode("utf-8", errors="ignore")
    finally:
        try:
            os.remove(temp_path)
        except OSError:
            pass


def extract_document_text(filename: str, contents: bytes) -> str:
    extension = get_document_extension(filename)
    if extension not in SUPPORTED_DOCUMENT_EXTENSIONS:
        raise ValueError("Unsupported file type. Please upload a PDF, PPT/PPTX, DOC, or DOCX file.")

    if extension == ".pdf":
        text = _extract_pdf_text(contents)
    elif extension == ".docx":
        text = _extract_docx_text(contents)
    elif extension in {".ppt", ".pptx"}:
        text = _extract_pptx_text(contents)
    else:
        text = _extract_legacy_doc_text(contents)

    normalized = _normalize_text(text)
    if not normalized:
        raise ValueError("No readable text could be extracted from the uploaded document.")

    return normalized


def build_document_topic(filename: str, fallback_topic: str | None = None) -> str:
    name = Path((filename or "").strip()).stem.replace("_", " ").replace("-", " ").strip()
    if name:
        return name[:100]
    fallback = (fallback_topic or "Uploaded Document").strip()
    return fallback[:100] or "Uploaded Document"